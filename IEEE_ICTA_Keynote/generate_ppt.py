"""
IEEE ICTA 2026 Keynote PPT Generator
From Cloud to Edge: Building AI Computing Infrastructure for the Age of Embodied Intelligence
27 slides, dark theme, 16:9
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from PIL import Image

# ── Image paths ──
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
IMG = {
    "circuit": os.path.join(IMG_DIR, "circuit_abstract.jpg"),
    "datacenter": os.path.join(IMG_DIR, "datacenter_corridor.jpg"),
    "server_rack": os.path.join(IMG_DIR, "server_rack.jpg"),
    "humanoid": os.path.join(IMG_DIR, "humanoid_robot.jpg"),
    "lingshen": os.path.join(IMG_DIR, "lingshen_1.jpg"),
    "tienkung": os.path.join(IMG_DIR, "tienkung_3.jpg"),
    "chip_wafer": os.path.join(IMG_DIR, "chip_wafer.jpg"),
}

# ── Colors ──
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SECTION = RGBColor(0x12, 0x12, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xA0)
LIGHT_GRAY = RGBColor(0xC0, 0xC0, 0xC0)
BLUE = RGBColor(0x0e, 0xa5, 0xe9)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
RED = RGBColor(0xef, 0x44, 0x44)
GOLD = RGBColor(0xff, 0xd7, 0x00)
DARK_BLUE = RGBColor(0x0a, 0x0a, 0x1a)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ──

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """lines = list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    for i, (text, fs, col, bld, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = col
        p.font.bold = bld
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(6)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_image_placeholder(slide, left, top, width, height, label="IMAGE"):
    """Add a dashed-border box as image placeholder"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(1)
    shape.line.dash_style = 4  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_image_fit(slide, img_path, left, top, width, height):
    """Add image scaled to fit within bounds, maintaining aspect ratio"""
    if not os.path.exists(img_path):
        add_image_placeholder(slide, left, top, width, height, f"Missing: {os.path.basename(img_path)}")
        return
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = width / height
    if img_ratio > box_ratio:
        # image is wider — fit to width
        new_w = width
        new_h = int(width / img_ratio)
        offset_top = top + (height - new_h) // 2
        offset_left = left
    else:
        # image is taller — fit to height
        new_h = height
        new_w = int(height * img_ratio)
        offset_left = left + (width - new_w) // 2
        offset_top = top
    slide.shapes.add_picture(img_path, offset_left, offset_top, new_w, new_h)

def add_fullbleed_image(slide, img_path, overlay_alpha=0.55):
    """Add full-bleed background image with dark overlay rectangle on top"""
    if not os.path.exists(img_path):
        add_image_placeholder(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, f"Missing: {os.path.basename(img_path)}")
        return
    # Add image covering full slide
    with Image.open(img_path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    slide_ratio = SLIDE_W / SLIDE_H
    if img_ratio > slide_ratio:
        # wider — fit height, crop width
        new_h = SLIDE_H
        new_w = int(SLIDE_H * img_ratio)
        left = -(new_w - SLIDE_W) // 2
        top = 0
    else:
        # taller — fit width, crop height
        new_w = SLIDE_W
        new_h = int(SLIDE_W / img_ratio)
        left = 0
        top = -(new_h - SLIDE_H) // 2
    slide.shapes.add_picture(img_path, left, top, new_w, new_h)
    # Dark overlay with transparency
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x15)
    # Set transparency via XML manipulation
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    sp_pr = overlay._element.find(f'.//{{{ns}}}solidFill')
    if sp_pr is not None:
        srgb = sp_pr.find(f'{{{ns}}}srgbClr')
        if srgb is not None:
            alpha_pct = int((1 - overlay_alpha) * 100000)
            alpha_el = etree.SubElement(srgb, f'{{{ns}}}alpha')
            alpha_el.set('val', str(alpha_pct))
    overlay.line.fill.background()

def set_notes(slide, text):
    """Set speaker notes for a slide"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_arrow_down(slide, cx, top, length, color=BLUE):
    """Draw a downward arrow using a triangle and line"""
    line_w = Pt(3)
    # vertical line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - Pt(1.5), top, Pt(3), length)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # arrow head
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - Inches(0.15), top + length, Inches(0.3), Inches(0.2))
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    tri.rotation = 180.0


# ════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# decorative bottom bar
add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), BLUE)

# Title
add_multi_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5), [
    ("From Cloud to Edge", 44, WHITE, True, PP_ALIGN.LEFT),
    ("Building AI Computing Infrastructure", 36, BLUE, False, PP_ALIGN.LEFT),
    ("for the Age of Embodied Intelligence", 36, BLUE, False, PP_ALIGN.LEFT),
])

# Author
add_multi_text(slide, Inches(1.5), Inches(5.0), Inches(8), Inches(2), [
    ("Chao Chen  |  \u9648\u8d85", 22, WHITE, True, PP_ALIGN.LEFT),
    ("General Manager, Suzhou Chuangjie Intelligent Technology", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Vice Chair, Supercomputing & Intelligent Computing Committee, CCIA", 14, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, GRAY, False, PP_ALIGN.LEFT),
    ("IEEE ICTA 2026  |  Singapore  |  October 2026", 14, ORANGE, False, PP_ALIGN.LEFT),
])

# decorative circuit image
add_image_fit(slide, IMG["circuit"], Inches(9.5), Inches(1.0), Inches(3.3), Inches(5.5))


# ════════════════════════════════════════════════════
# SLIDE 2 — About Me
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "About Me", 36, WHITE, True)

# Three columns timeline
cols = [
    ("\U0001f4d0 Quantitative Trading", BLUE,
     ["Applied Mathematics", "High-frequency computation", "", '"Compute = Profit"']),
    ("\U0001f3d7\ufe0f Traditional AI Infra", ORANGE,
     ["Operating GPU clusters", "at national scale (DC-centric)", "", '"Compute = Infrastructure"']),
    ("\U0001f916 Embodied Intelligence", GREEN,
     ["Investing in humanoid robots", "Compute leaves the DC", "", '"Compute = Autonomy"']),
]

for i, (title, color, bullets) in enumerate(cols):
    x = Inches(0.8 + i * 4.0)
    # colored top bar
    add_rect(slide, x, Inches(2.0), Inches(3.5), Inches(0.08), color)
    # title
    add_textbox(slide, x, Inches(2.2), Inches(3.5), Inches(0.7), title, 22, color, True)
    # bullets
    lines = [(b, 16, LIGHT_GRAY if b and not b.startswith('"') else GRAY, b.startswith('"'), PP_ALIGN.LEFT) for b in bullets]
    add_multi_text(slide, x, Inches(3.0), Inches(3.5), Inches(3.0), lines)

# Timeline arrows between columns
for i in range(2):
    x = Inches(4.3 + i * 4.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.3), Inches(0.5), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()


# ════════════════════════════════════════════════════
# SLIDE 3 — Three Waves of Compute
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Three Waves of Compute Demand", 36, WHITE, True)

# Wave cards — Wave 2 = traditional AI (stays in DC), Wave 3 = embodied AI (walks out)
waves = [
    ("Wave 1", "Quantitative Finance (2010s)", BLUE,
     ["\u2022 Microsecond latency matters", "\u2022 Every compute cycle = money",
      "", "\u25b8 Compute = Profit"]),
    ("Wave 2", "Traditional AI Training (2020s)", ORANGE,
     ["\u2022 GPU clusters \u2014 stays in the data center", "\u2022 Petaflops throughput, 700W/chip",
      "\u2022 Liquid cooling, MW-scale power", "", "\u25b8 Compute = Infrastructure"]),
    ("Wave 3", "Embodied AI (2026\u2192)", GREEN,
     ["\u2022 Compute leaves the data center", "\u2022 On-body inference, <50W budget",
      "\u2022 Real-time safety-critical control", "", "\u25b8 Compute = Autonomy"]),
]

for i, (wave, subtitle, color, bullets) in enumerate(waves):
    x = Inches(0.8 + i * 4.0)
    # card background
    add_rounded_rect(slide, x, Inches(1.8), Inches(3.6), Inches(4.2),
                     RGBColor(0x20, 0x20, 0x38), color)
    add_textbox(slide, x + Inches(0.3), Inches(2.0), Inches(3.0), Inches(0.6),
                wave, 28, color, True)
    add_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(3.0), Inches(0.6),
                subtitle, 16, WHITE, False)
    lines = [(b, 15, GOLD if b.startswith("\u25b8") else LIGHT_GRAY, b.startswith("\u25b8"), PP_ALIGN.LEFT) for b in bullets]
    add_multi_text(slide, x + Inches(0.3), Inches(3.3), Inches(3.0), Inches(2.5), lines)

# Bottom quote
add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            '"Each wave demands fundamentally different IC design."', 20, GOLD, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 4 — The Numbers
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "AI Compute: The Numbers", 36, WHITE, True)

numbers = [
    ("2/3", "of all AI compute\nis inference (2026)", BLUE),
    ("$50B+", "inference chip\nmarket (2026)", ORANGE),
    ("31%", "embodied AI\nmarket CAGR", GREEN),
]

for i, (num, desc, color) in enumerate(numbers):
    x = Inches(1.0 + i * 4.0)
    add_textbox(slide, x, Inches(2.2), Inches(3.5), Inches(2.0),
                num, 72, color, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(4.5), Inches(3.5), Inches(1.5),
                desc, 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Source
add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
            "Source: Deloitte TMT Predictions 2026  |  Omdia Market Radar 2026", 12, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 5 — Section Divider: Cloud
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)

add_fullbleed_image(slide, IMG["datacenter"], overlay_alpha=0.55)

add_multi_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4.0), [
    ("PART I", 20, BLUE, True, PP_ALIGN.LEFT),
    ("The Cloud", 56, WHITE, True, PP_ALIGN.LEFT),
    ("Operating AI at Scale", 28, LIGHT_GRAY, False, PP_ALIGN.LEFT),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 6 — Scale of the Challenge
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "China's Intelligent Computing Centers", 36, WHITE, True)

stats = [
    ("200+", "intelligent computing centers built since 2023", RED),
    ("60,000+", "GPU accelerators deployed at leading operators", ORANGE),
    ("12+ EFlops", "total deployed AI compute (and growing)", BLUE),
    ("<30%", "average GPU utilization rate", WHITE),
]

for i, (num, desc, color) in enumerate(stats):
    y = Inches(1.8 + i * 1.2)
    add_textbox(slide, Inches(1.5), y, Inches(3), Inches(1), num, 48, color, True)
    add_textbox(slide, Inches(5), y + Inches(0.15), Inches(7), Inches(0.8), desc, 20, LIGHT_GRAY, False)

add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.8),
            '"The system-level challenges dominate everything."', 20, GOLD, True)


# ════════════════════════════════════════════════════
# SLIDE 7 — Challenge 1: Thermal
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Challenge #1: Thermal Management", 36, WHITE, True)

# Left column
add_multi_text(slide, Inches(1), Inches(1.8), Inches(5), Inches(3), [
    ("Single accelerator:", 18, GRAY, False, PP_ALIGN.LEFT),
    ("700W+", 44, RED, True, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Single rack:", 18, GRAY, False, PP_ALIGN.LEFT),
    ("40\u201380 kW", 44, ORANGE, True, PP_ALIGN.LEFT),
])

# Right column
add_multi_text(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3.5), [
    ("Air Cooling \u2192 Liquid Cooling", 24, BLUE, True, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("PUE: 1.4 \u2192 1.15", 28, GREEN, True, PP_ALIGN.LEFT),
    ("= millions saved per year", 18, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

add_image_fit(slide, IMG["server_rack"], Inches(7), Inches(4.2), Inches(5), Inches(2.5))

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            '"Your chip\'s TDP is not a spec \u2014 it\'s a cost multiplier."', 20, GOLD, True)


# ════════════════════════════════════════════════════
# SLIDE 8 — Challenge 2: Power
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Challenge #2: Power Supply", 36, WHITE, True)

add_multi_text(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(2.5), [
    ("1,000 GPUs = 1 MW", 28, ORANGE, True, PP_ALIGN.LEFT),
    ("10,000 GPUs = 10 MW", 28, RED, True, PP_ALIGN.LEFT),
    ("= a small town's power consumption", 18, GRAY, False, PP_ALIGN.LEFT),
])

# East vs West comparison
add_rounded_rect(slide, Inches(1), Inches(4.0), Inches(5), Inches(1.2),
                 RGBColor(0x20, 0x20, 0x38), RED)
add_textbox(slide, Inches(1.3), Inches(4.1), Inches(4.5), Inches(1.0),
            "East China:  \u00a50.8/kWh  \u2192  Expensive", 20, RED, True)

add_rounded_rect(slide, Inches(1), Inches(5.4), Inches(5), Inches(1.2),
                 RGBColor(0x20, 0x20, 0x38), GREEN)
add_textbox(slide, Inches(1.3), Inches(5.5), Inches(4.5), Inches(1.0),
            "West China:  \u00a50.3/kWh  \u2192  AI moves west", 20, GREEN, True)

# Visual: East vs West infographic (replaces map placeholder)
# "East Data, West Computing" visual
add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.8),
                 RGBColor(0x18, 0x18, 0x30), RGBColor(0x30, 0x30, 0x50))

# Header
add_textbox(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
            "\u201cEastern Data, Western Computing\u201d", 18, GOLD, True, PP_ALIGN.CENTER)

# West side (green = cheap)
add_rounded_rect(slide, Inches(7.3), Inches(2.6), Inches(2.3), Inches(3.5),
                 RGBColor(0x15, 0x30, 0x20), GREEN)
add_multi_text(slide, Inches(7.4), Inches(2.7), Inches(2.1), Inches(3.3), [
    ("WEST", 16, GREEN, True, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("Guizhou \u00b7 Inner Mongolia", 13, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Gansu \u00b7 Xinjiang", 13, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("\u00a50.3/kWh", 24, GREEN, True, PP_ALIGN.CENTER),
    ("Abundant land", 12, GRAY, False, PP_ALIGN.CENTER),
    ("Clean energy + waste", 12, GRAY, False, PP_ALIGN.CENTER),
    ("heat recovery", 12, GRAY, False, PP_ALIGN.CENTER),
])

# Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(9.7), Inches(3.8), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = GOLD
arrow.line.fill.background()

# East side (red = expensive)
add_rounded_rect(slide, Inches(10.6), Inches(2.6), Inches(1.7), Inches(3.5),
                 RGBColor(0x30, 0x15, 0x15), RED)
add_multi_text(slide, Inches(10.7), Inches(2.7), Inches(1.5), Inches(3.3), [
    ("EAST", 16, RED, True, PP_ALIGN.CENTER),
    ("", 6, WHITE, False, PP_ALIGN.CENTER),
    ("Beijing", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Shanghai", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Shenzhen", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 8, WHITE, False, PP_ALIGN.CENTER),
    ("\u00a50.8/kWh", 24, RED, True, PP_ALIGN.CENTER),
    ("Limited land", 12, GRAY, False, PP_ALIGN.CENTER),
    ("High demand", 12, GRAY, False, PP_ALIGN.CENTER),
])

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            '"Power is the new oil of the AI era."', 20, GOLD, True)


# ════════════════════════════════════════════════════
# SLIDE 9 — Challenge 3: Utilization
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Challenge #3: GPU Utilization", 36, WHITE, True)

# Gauge visual - big number
add_textbox(slide, Inches(1), Inches(1.8), Inches(5), Inches(1.5),
            "<30%", 80, RED, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.5), Inches(5), Inches(0.6),
            "Average GPU Utilization", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Why list
reasons = [
    "Multi-tenant scheduling is hard",
    "Memory fragmentation between model sizes",
    "Daily hardware failures at scale",
    "Recovery and job migration overhead",
]
add_textbox(slide, Inches(7), Inches(1.8), Inches(5), Inches(0.6), "Why?", 28, ORANGE, True)
lines = [("\u2022 " + r, 18, LIGHT_GRAY, False, PP_ALIGN.LEFT) for r in reasons]
add_multi_text(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(2.5), lines)

# Solution hint
add_rounded_rect(slide, Inches(7), Inches(5.2), Inches(5.5), Inches(0.9),
                 RGBColor(0x15, 0x25, 0x15), GREEN)
add_textbox(slide, Inches(7.2), Inches(5.25), Inches(5.1), Inches(0.8),
            "\u2192 Heterogeneous pooling + vGPU slicing can lift to 70%+", 16, GREEN, True)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            '"Utilization rate matters more than total capacity."', 20, GOLD, True)


# ════════════════════════════════════════════════════
# SLIDE 10 — Cloud Lessons for IC
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
            "What Cloud Operations Teach IC Designers", 36, WHITE, True)

lessons = [
    "Design for system deployability, not benchmarks",
    "Performance per watt > peak performance",
    "Build diagnostics into the silicon",
    "Reliability over years, not hours",
]

for i, lesson in enumerate(lessons):
    y = Inches(1.8 + i * 1.2)
    # green bar
    add_rect(slide, Inches(1.3), y, Inches(0.08), Inches(0.8), GREEN)
    # check + text
    add_textbox(slide, Inches(1.7), y, Inches(1), Inches(0.8), "\u2705", 24, GREEN, True)
    add_textbox(slide, Inches(2.3), y + Inches(0.1), Inches(9), Inches(0.7),
                f"{i+1}. {lesson}", 24, WHITE, False)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
            '"The best chip is the one that runs stable at 85% load for 365 days."', 18, GOLD, True)


# ════════════════════════════════════════════════════
# SLIDE 11 — Section Divider: Embodied AI
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)

add_fullbleed_image(slide, IMG["humanoid"], overlay_alpha=0.55)

add_multi_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4.0), [
    ("PART II", 20, GREEN, True, PP_ALIGN.LEFT),
    ("From Racks to Bodies", 52, WHITE, True, PP_ALIGN.LEFT),
    ("When Compute Walks Out of the Data Center", 24, LIGHT_GRAY, False, PP_ALIGN.LEFT),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 12 — Why Embodied Intelligence
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Why I Invest in Embodied AI", 36, WHITE, True)

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
            '"I asked the engineering teams:\n What is your biggest bottleneck?"', 24, LIGHT_GRAY, True,
            PP_ALIGN.LEFT)

answers = [
    ("\u274c  Not funding", RED),
    ("\u274c  Not algorithms", RED),
    ("\u2705  On-body compute", GREEN),
]
for i, (text, color) in enumerate(answers):
    y = Inches(3.2 + i * 0.9)
    add_textbox(slide, Inches(2), y, Inches(8), Inches(0.8), text, 28, color, True)

add_multi_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.5), [
    ('"The available processors simply cannot deliver', 20, GOLD, True, PP_ALIGN.LEFT),
    (' what we need within our power and size constraints."', 20, GOLD, True, PP_ALIGN.LEFT),
])


# ════════════════════════════════════════════════════
# SLIDE 13 — LingShen Technology
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Portfolio: LingShen Technology  \u7075\u751f\u79d1\u6280", 34, WHITE, True)

# green accent bar
add_rect(slide, Inches(1), Inches(1.5), Inches(11.5), Inches(0.06), GREEN)

bullets = [
    "\u2022  Tsinghua University origin, founded 2023",
    "\u2022  \u201cUniversal Brain\u201d platform for humanoid robots",
    "\u2022  4-in-1 embodied large model:",
    "     \u25b8  Perception + Movement + Manipulation + Planning",
    "\u2022  L1 Robot: 26 DOF, 85 kg, NVIDIA Orin, 12h+ endurance",
    "\u2022  <300 samples for new scenario, deploy in 5 hours",
    "\u2022  Deployed: smart manufacturing, hospital rounds, DC inspection",
    "\u2022  Pre-A / Pre-A+ funding: >100M RMB",
]
lines = [(b, 18, LIGHT_GRAY, False, PP_ALIGN.LEFT) for b in bullets]
add_multi_text(slide, Inches(1.2), Inches(1.8), Inches(6.5), Inches(5), lines)

add_image_fit(slide, IMG["lingshen"], Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.5))


# ════════════════════════════════════════════════════
# SLIDE 14 — TIEN Kung
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Portfolio: TIEN Kung  \u5929\u5de5", 34, WHITE, True)

add_rect(slide, Inches(1), Inches(1.5), Inches(11.5), Inches(0.06), ORANGE)

bullets = [
    "\u2022  Beijing Innovation Center of Humanoid Robotics (MIIT)",
    "\u2022  World\u2019s first pure-electric full-size humanoid running at 6 km/h",
    "\u2022  v2.0 Pro: 173 cm, 73 kg, 42 DOF (incl. dexterous hands)",
    "\u2022  On-body compute: dual NVIDIA Orin AGX \u2192 550 TOPS",
    "\u2022  Payload: 4 kg/arm (16 kg bimanual), 3.5 h endurance",
    "\u2022  Open-source universal robot mother platform",
]
lines = [(b, 18, LIGHT_GRAY, False, PP_ALIGN.LEFT) for b in bullets]
add_multi_text(slide, Inches(1.2), Inches(1.8), Inches(6.5), Inches(4.5), lines)

add_image_fit(slide, IMG["tienkung"], Inches(8.2), Inches(1.8), Inches(4.5), Inches(4.5))


# ════════════════════════════════════════════════════
# SLIDE 15 — Real-World Deployments (NEW)
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.7),
            "Where Embodied AI Is Already Working", 34, WHITE, True)

add_rect(slide, Inches(1), Inches(1.2), Inches(11.5), Inches(0.06), GREEN)

# Three deployment case columns
cases = [
    ("Smart Manufacturing", GREEN,
     ["\u2022 Injection molding + stamping",
      "  production lines",
      "\u2022 Grasping, sorting, quality",
      "  inspection, assembly",
      "\u2022 ROI: replaces 2\u20133 manual",
      "  workers per robot",
      "",
      "LingShen L1 deployed"]),
    ("Hospital Ward Rounds", BLUE,
     ["\u2022 Automated vital signs",
      "  (BP, SpO\u2082, temperature)",
      "\u2022 30 wards/morning, data",
      "  synced to HIS system",
      "\u2022 Medical LLM + voice",
      "  interaction for pre-diagnosis",
      "",
      "LingShen L1 deployed"]),
    ("Data Center Inspection", ORANGE,
     ["\u2022 24/7 autonomous patrol",
      "  of server racks",
      "\u2022 Anomaly detection: power-on/",
      "  off, restart, LED status",
      "\u2022 <300 samples to train",
      "  for new DC environment",
      "",
      "LingShen L1 deployed"]),
]

for i, (title, color, bullets) in enumerate(cases):
    x = Inches(0.8 + i * 4.2)
    add_rounded_rect(slide, x, Inches(1.5), Inches(3.8), Inches(5.2),
                     RGBColor(0x1a, 0x1a, 0x30), color)
    add_textbox(slide, x + Inches(0.2), Inches(1.6), Inches(3.4), Inches(0.5),
                title, 20, color, True)
    add_rect(slide, x + Inches(0.2), Inches(2.15), Inches(3.4), Inches(0.04), color)
    lines = [(b, 14, GOLD if "deployed" in b.lower() else LIGHT_GRAY,
              "deployed" in b.lower(), PP_ALIGN.LEFT) for b in bullets]
    add_multi_text(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(4.2), lines)

add_textbox(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.5),
            "TIEN Kung 2.0: open-source platform enabling logistics, elder care, hazardous environments", 14, GRAY, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 16 — Cloud vs On-Body Comparison (was 15)
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
            "Cloud Chips vs. On-Body Chips: Two Worlds", 34, WHITE, True)

# Table
rows = [
    ("", "Cloud", "On-Body"),
    ("Compute", "PetaFLOPS", "550 TOPS (dual Orin)"),
    ("Power", "700W (liquid cooled)", "<50W (battery)"),
    ("Latency", "~1 second OK", "<10ms required"),
    ("Failure mode", "Swap & restart", "Cannot fail"),
    ("Workload", "Single model", "Vision+LLM+Motor+Tactile"),
    ("Form factor", "Rack-mounted", "Fits in 73 kg torso"),
    ("Cost target", "$30,000/card", "<$500/unit"),
]

for i, row in enumerate(rows):
    y = Inches(1.4 + i * 0.7)
    bg_color = RGBColor(0x25, 0x25, 0x40) if i % 2 == 0 else BG_DARK
    if i == 0:
        bg_color = RGBColor(0x15, 0x15, 0x30)

    add_rect(slide, Inches(0.8), y, Inches(3.5), Inches(0.65), bg_color)
    add_rect(slide, Inches(4.35), y, Inches(4.2), Inches(0.65), bg_color)
    add_rect(slide, Inches(8.6), y, Inches(4.2), Inches(0.65), bg_color)

    label_color = GRAY if i == 0 else LIGHT_GRAY
    cloud_color = BLUE if i == 0 else WHITE
    body_color = GREEN if i == 0 else WHITE

    add_textbox(slide, Inches(1.0), y + Inches(0.05), Inches(3.2), Inches(0.55),
                row[0], 16 if i == 0 else 18, label_color, i == 0, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(4.55), y + Inches(0.05), Inches(3.8), Inches(0.55),
                row[1], 16 if i == 0 else 18, cloud_color, i == 0, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.8), y + Inches(0.05), Inches(3.8), Inches(0.55),
                row[2], 16 if i == 0 else 18, body_color, i == 0, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 16 — The Core Insight
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_multi_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(5), [
    ('"Cloud computing\'s challenge:', 30, WHITE, True, PP_ALIGN.CENTER),
    (' fitting chips into racks.', 30, BLUE, True, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    (' Embodied computing\'s challenge:', 30, WHITE, True, PP_ALIGN.CENTER),
    (' fitting racks into bodies."', 30, GREEN, True, PP_ALIGN.CENTER),
    ("", 30, WHITE, False, PP_ALIGN.CENTER),
    ("These require fundamentally different IC design philosophies.", 18, GRAY, False, PP_ALIGN.CENTER),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 17 — Three-Tier Architecture
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
            "The Three-Tier Compute Architecture for Embodied AI", 32, WHITE, True)

tiers = [
    ("\u2601\ufe0f  CLOUD  \u201cBig Brain\u201d", "GPU Cluster  |  Heterogeneous pool  |  EFlops",
     "Model training, world model updates, data flywheel", "Hours / Days", BLUE, Inches(1.2)),
    ("\U0001f3ed  EDGE  \u201cMid Brain\u201d", "Edge Server  |  100\u2013300W  |  Multi-robot coordination",
     "Scene inference, swarm orchestration, 5G/6G uplink", "~100ms", ORANGE, Inches(3.2)),
    ("\U0001f916  ON-BODY  \u201cSmall Brain\u201d", "SoC / NPU  |  <50W  |  550 TOPS",
     "Perception, safety-critical control, autonomous action", "<10ms", GREEN, Inches(5.2)),
]

for emoji_title, subtitle, desc, latency, color, y in tiers:
    add_rounded_rect(slide, Inches(2.5), y, Inches(8.5), Inches(1.6),
                     RGBColor(0x20, 0x20, 0x38), color)
    add_textbox(slide, Inches(2.8), y + Inches(0.1), Inches(4), Inches(0.5),
                emoji_title, 22, color, True)
    add_textbox(slide, Inches(2.8), y + Inches(0.55), Inches(4), Inches(0.4),
                subtitle, 14, GRAY, False)
    add_textbox(slide, Inches(7), y + Inches(0.15), Inches(3.5), Inches(0.5),
                desc, 14, LIGHT_GRAY, False)
    add_textbox(slide, Inches(7), y + Inches(0.65), Inches(3.5), Inches(0.4),
                f"Latency: {latency}", 13, color, True)

# Arrows between tiers
for y_start in [Inches(2.8), Inches(4.8)]:
    cx = Inches(6.75)
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.15), y_start, Inches(0.3), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()


# ════════════════════════════════════════════════════
# SLIDE 18 — IC Demands per Tier
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "IC Demands by Compute Tier", 36, WHITE, True)

headers = ["", "Cloud", "Edge", "On-Body"]
data_rows = [
    ("Key IC type", "GPU / AI\naccelerator", "Inference\naccelerator", "Heterogeneous\nSoC + NPU"),
    ("Power", "400\u2013700W", "100\u2013300W", "<50W (battery)"),
    ("Real example", "60K+ GPUs\nin pooled clusters", "5G/6G connected\nedge servers", "550 TOPS\n(dual Orin AGX)"),
    ("Key challenge", "Interconnect\n& utilization", "Multi-model\nscheduling", "Multi-modal\nintegration"),
    ("Market status", "Mature", "Exploding", "GREENFIELD\nBIGGEST OPPORTUNITY"),
]

col_colors = [GRAY, BLUE, ORANGE, GREEN]
col_x = [Inches(0.8), Inches(3.8), Inches(6.6), Inches(9.4)]
col_w = [Inches(2.8), Inches(2.6), Inches(2.6), Inches(3.2)]

for j, (header, color) in enumerate(zip(headers, col_colors)):
    add_textbox(slide, col_x[j], Inches(1.6), col_w[j], Inches(0.5),
                header, 18, color, True, PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

for i, (label, *vals) in enumerate(data_rows):
    y = Inches(2.2 + i * 1.0)
    bg = RGBColor(0x20, 0x20, 0x35) if i % 2 == 0 else BG_DARK
    add_rect(slide, Inches(0.6), y, Inches(12.2), Inches(0.95), bg)
    add_textbox(slide, col_x[0], y + Inches(0.05), col_w[0], Inches(0.85),
                label, 15, GRAY, True)
    for j, val in enumerate(vals):
        color = GREEN if "BIGGEST" in val else WHITE
        add_textbox(slide, col_x[j+1], y + Inches(0.05), col_w[j+1], Inches(0.85),
                    val, 15, color, "BIGGEST" in val, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 19 — Section Divider: Recommendations
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_SECTION)

add_fullbleed_image(slide, IMG["chip_wafer"], overlay_alpha=0.55)

add_multi_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4.0), [
    ("PART III", 20, ORANGE, True, PP_ALIGN.LEFT),
    ("Recommendations", 52, WHITE, True, PP_ALIGN.LEFT),
    ("for IC Designers", 28, LIGHT_GRAY, False, PP_ALIGN.LEFT),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 20 — For Cloud Chip Designers
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
            "For Cloud / Data Center Chip Designers", 34, WHITE, True)

# Blue left bar
add_rect(slide, Inches(1.2), Inches(1.5), Inches(0.06), Inches(5.2), BLUE)

recs = [
    ("Performance per watt > peak performance", "Your customer's #1 cost is electricity"),
    ("Standardize interconnects", "Proprietary = lock-in = slow adoption"),
    ("Build diagnostics into silicon", "Remote monitoring is an operational necessity"),
    ("Design for multi-tenant workloads", "Hardware isolation improves utilization"),
]

for i, (title, desc) in enumerate(recs):
    y = Inches(1.7 + i * 1.25)
    add_textbox(slide, Inches(1.8), y, Inches(10), Inches(0.5),
                f"{i+1}.  {title}", 24, WHITE, True)
    add_textbox(slide, Inches(2.3), y + Inches(0.55), Inches(9), Inches(0.5),
                desc, 16, GRAY, False)


# ════════════════════════════════════════════════════
# SLIDE 21 — For Embodied Chip Designers
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
            "For Edge / Embodied Chip Designers", 34, WHITE, True)

add_rect(slide, Inches(1.2), Inches(1.5), Inches(0.06), Inches(5.5), GREEN)

recs = [
    ("Heterogeneous compute is mandatory", "CNN + Transformer + Control on one die", WHITE),
    ("Power is the hard constraint", "Design to 50W first, then maximize compute", WHITE),
    ("Latency > throughput", "One result in 5ms, not 1000 results per second", WHITE),
    ("Safety is non-negotiable", "Functional safety, error correction, graceful degradation", RED),
    ("Cost at scale: <$500/unit", 'A $200K robot will never reach mass adoption', WHITE),
]

for i, (title, desc, highlight) in enumerate(recs):
    y = Inches(1.6 + i * 1.05)
    add_textbox(slide, Inches(1.8), y, Inches(10), Inches(0.45),
                f"{i+1}.  {title}", 22, highlight, True)
    add_textbox(slide, Inches(2.3), y + Inches(0.45), Inches(9), Inches(0.45),
                desc, 15, GRAY, False)


# ════════════════════════════════════════════════════
# SLIDE 22 — The Bridging Insight
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_multi_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(5.5), [
    ('"The best chip is never the one', 28, WHITE, True, PP_ALIGN.CENTER),
    (' with the highest benchmark score.', 28, WHITE, True, PP_ALIGN.CENTER),
    ("", 20, WHITE, False, PP_ALIGN.CENTER),
    (" It is the one that runs reliably,", 28, BLUE, True, PP_ALIGN.CENTER),
    (" efficiently, and affordably", 28, BLUE, True, PP_ALIGN.CENTER),
    (' in its target deployment environment."', 28, BLUE, True, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    ("\u2014 From quantitative trading,", 20, GRAY, False, PP_ALIGN.CENTER),
    ("  through data centers, to robots.", 20, GRAY, False, PP_ALIGN.CENTER),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 23 — Looking Ahead
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "Looking Ahead", 36, WHITE, True)

timeline = [
    ("2026", '"Inference Famine"', "Demand for inference compute outpaces supply", BLUE),
    ("2027", "First wave of mass-produced humanoid robots", "Logistics, manufacturing, elder care", ORANGE),
    ("2028+", "On-body AI compute market", "Comparable in scale to smartphone SoC market", GREEN),
]

for i, (year, title, desc, color) in enumerate(timeline):
    y = Inches(1.8 + i * 1.7)
    # Year circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    p = tf.paragraphs[0]
    p.text = year
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Vertical line between circles
    if i < 2:
        add_rect(slide, Inches(2.07), y + Inches(1.2), Inches(0.06), Inches(0.5), color)

    # Text
    add_textbox(slide, Inches(3.2), y + Inches(0.05), Inches(8), Inches(0.6),
                title, 24, WHITE, True)
    add_textbox(slide, Inches(3.2), y + Inches(0.6), Inches(8), Inches(0.5),
                desc, 16, GRAY, False)

# Bottom
add_multi_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(1.0), [
    ("Every robot = a mobile computing platform", 18, GOLD, True, PP_ALIGN.CENTER),
    ("Aggregate demand = enormous new IC market", 18, GOLD, True, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════════════════
# SLIDE 24 — Convergence Visual
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "The Convergence", 36, WHITE, True)

# Three overlapping circles (Venn diagram effect)
circles_data = [
    (Inches(4.0), Inches(1.8), BLUE, "Cloud AI\nTraining chips\n(mature)"),
    (Inches(6.5), Inches(1.8), ORANGE, "Edge AI\nInference chips\n(growing)"),
    (Inches(5.25), Inches(3.8), GREEN, "On-Body AI\nEmbodied SoCs\n(emerging)"),
]

for cx, cy, color, label in circles_data:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(3.5), Inches(3.5))
    circle.fill.solid()
    # Semi-transparent effect via darker fill
    circle.fill.fore_color.rgb = RGBColor(color[0] // 3, color[1] // 3, color[2] // 3)
    circle.line.color.rgb = color
    circle.line.width = Pt(2)
    tf = circle.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except:
        pass
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Center label
add_rounded_rect(slide, Inches(5.4), Inches(3.8), Inches(2.8), Inches(0.7),
                 RGBColor(0x30, 0x30, 0x50), GOLD)
add_textbox(slide, Inches(5.4), Inches(3.85), Inches(2.8), Inches(0.6),
            "The IC Opportunity", 16, GOLD, True, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 25 — Call to Action
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
            "For Everyone in This Room", 36, WHITE, True)

add_multi_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5), [
    ("You are designing the", 28, WHITE, False, PP_ALIGN.LEFT),
    ("nervous system", 36, GOLD, True, PP_ALIGN.LEFT),
    ("of tomorrow's intelligent machines.", 28, WHITE, False, PP_ALIGN.LEFT),
    ("", 20, WHITE, False, PP_ALIGN.LEFT),
    ("Your decisions on:", 24, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("    \u2022  Architecture", 24, BLUE, True, PP_ALIGN.LEFT),
    ("    \u2022  Power efficiency", 24, ORANGE, True, PP_ALIGN.LEFT),
    ("    \u2022  Reliability", 24, GREEN, True, PP_ALIGN.LEFT),
    ("", 16, WHITE, False, PP_ALIGN.LEFT),
    ("will determine whether embodied AI", 24, WHITE, False, PP_ALIGN.LEFT),
    ("becomes reality \u2014 or stays in the lab.", 24, WHITE, True, PP_ALIGN.LEFT),
])


# ════════════════════════════════════════════════════
# SLIDE 26 — Closing Quote
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_multi_text(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(5.5), [
    ('"I started optimizing compute for financial models,', 24, WHITE, False, PP_ALIGN.CENTER),
    (' where microseconds meant money.', 24, BLUE, True, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    (" I then built computing centers,", 24, WHITE, False, PP_ALIGN.CENTER),
    (" where power efficiency meant survival.", 24, ORANGE, True, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    (" Now I invest in robots,", 24, WHITE, False, PP_ALIGN.CENTER),
    (' where real-time inference means safety."', 24, GREEN, True, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    ("Advancing IC with AI.", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Advancing AI with IC \u2014", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("from the cloud, through the edge, into the body.", 22, GOLD, True, PP_ALIGN.CENTER),
], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════
# SLIDE 27 — Thank You
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# decorative top bar
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.3), BLUE)

add_multi_text(slide, Inches(2), Inches(1.5), Inches(9), Inches(5), [
    ("Thank You", 56, WHITE, True, PP_ALIGN.CENTER),
    ("\u8c22\u8c22", 36, GRAY, False, PP_ALIGN.CENTER),
    ("", 30, WHITE, False, PP_ALIGN.CENTER),
    ("Chao Chen  \u9648\u8d85", 28, WHITE, True, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.CENTER),
    ("Suzhou Chuangjie Intelligent Technology Co., Ltd.", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("\u82cf\u5dde\u521b\u754c\u667a\u80fd\u79d1\u6280\u6709\u9650\u516c\u53f8", 16, GRAY, False, PP_ALIGN.CENTER),
    ("", 16, WHITE, False, PP_ALIGN.CENTER),
    ("IEEE ICTA 2026  |  Singapore", 18, ORANGE, False, PP_ALIGN.CENTER),
], anchor=MSO_ANCHOR.MIDDLE)

# decorative bottom bar
add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), BLUE)


# ── Speaker Notes ──
NOTES = {
    0: """Good morning, everyone. Thank you to the IEEE ICTA organizing committee for the invitation. It's a privilege to be here at a conference dedicated to advancing integrated circuits with AI.

My name is Chao Chen. I'm going to talk about AI computing infrastructure — but not from the chip designer's perspective. I'm going to talk about it from the other side: the perspective of someone who deploys chips at scale, who operates the data centers they go into, and who invests in the robots that will need them next.""",

    1: """Let me start with a personal observation. In my career, I've lived through what I call three waves of compute demand.

The first wave was quantitative trading. I started in quantitative finance, building systems where microseconds mattered and computational throughput directly translated to profit. Computing power is not an IT cost — it's the core production asset.

The second wave — what I call "traditional AI" — is large-scale GPU clusters for training foundation models. This wave is data-center-centric: the chips stay in the rack, the cooling stays in the building, and the power stays on the grid. All the hard problems — thermal management, power delivery, interconnection, utilization — are system-level problems, but they happen in a controlled environment. 700 watts per chip is fine when you have liquid cooling. Megawatts of power is fine when you build next to a hydroelectric dam.

The third wave is fundamentally different: embodied intelligence. Compute leaves the data center. It walks out on two legs, runs on batteries, and must make safety-critical decisions in milliseconds. I'm an investor in LingShen Technology and TIEN Kung, and through these investments I've seen computing demands that break every assumption from the data center world.""",

    2: """According to Deloitte, inference workloads will account for roughly two-thirds of all AI compute in 2026 — up from one-third in 2023. The market for inference-optimized chips alone will exceed 50 billion US dollars this year. And the embodied AI market is projected to grow at a 31% compound annual growth rate.

The demand for compute is not just growing — it's diversifying. Each new form of compute creates new demands on IC design.

Today I want to take you on a journey from the cloud to the edge to the robot's body — and at each stop, I'll share what chip designers need to know.""",

    3: """[Data slide — let the numbers speak]

Two-thirds of all AI compute is inference. 50 billion dollar inference chip market. 31% CAGR for embodied AI.

Source: Deloitte TMT Predictions 2026, Omdia Market Radar 2026.""",

    4: """Let me start with the data center.

[Section transition — pause briefly]""",

    5: """China has built over 200 intelligent computing centers in the past three years. I've been directly involved in several of these. The largest operators now manage over 60,000 GPU accelerators across multiple sites, with total deployed compute exceeding 12 exaFLOPS — spread across centers in Yichang, Ordos, Sichuan, Xinjiang, and more.

To give you a sense of scale: a single facility in Ordos deployed 3,000 petaFLOPS with a 10-billion-yuan investment. The largest planned facility — in Wenzhou — is targeting 50,000 petaFLOPS with investment of nearly 14 billion yuan. These are massive infrastructure projects comparable to building power plants.

Yet even at this scale, the most important lessons I've learned have nothing to do with chip specifications. When you operate thousands of GPUs in a single facility, the system-level challenges dominate everything. Let me walk you through the top three.""",

    6: """Challenge number one: thermal management.

A single high-end AI accelerator now consumes 700 watts or more. A fully loaded rack can draw 40 to 80 kilowatts. At this power density, traditional air cooling simply cannot keep up.

In our facilities, we've transitioned from air cooling to liquid cooling, bringing our PUE from approximately 1.4 down to 1.15. That 0.25 reduction translates directly to millions of dollars in annual energy savings and, more importantly, to stable chip operating temperatures that extend hardware lifetime.

The takeaway for IC designers: your chip's thermal design power is not just a spec on a datasheet. It directly determines what kind of cooling infrastructure is required, which determines the total cost of ownership. Reducing TDP by even 50 watts can change the economics of an entire data center deployment.""",

    7: """Challenge number two: power supply.

One thousand GPU accelerators consume approximately one megawatt. A large center may have ten to fifty thousand cards — 10 to 50 megawatts — the electrical load of a small town. The largest planned facilities now reach 50,000 petaFLOPS, with investment exceeding one billion US dollars.

In China, this has led to a geographic redistribution of AI infrastructure — the "Eastern Data, Western Computing" national strategy. New centers are being built in Guizhou, Inner Mongolia, Gansu, and Xinjiang, where electricity costs one-third of eastern prices and renewable energy is abundant.

An interesting innovation: some operators are recovering waste heat from their cooling systems to power agricultural operations — aquaculture and hydroponic farming — turning an energy cost into a revenue stream. This is the kind of system-level thinking that matters.

For chip designers: performance per watt is arguably more important than absolute performance. A chip that delivers 90% of the FLOPS at 70% of the power will win in large-scale deployment every time.""",

    8: """Challenge number three: utilization.

This one surprised me. The average GPU utilization rate across China's intelligent computing centers is estimated at less than 30 percent. We buy enormous compute capacity, and most of it sits idle.

Why? Multi-tenant scheduling is hard. Memory fragmentation between models of different sizes is a real problem. When a card fails — which happens daily at scale — recovery wastes hours of cluster time.

The good news is that heterogeneous resource pooling and vGPU slicing techniques are emerging — some operators have demonstrated that these can lift utilization from under 30% to over 70%, reducing cost by 40-70% compared to traditional cloud. But the fundamental point stands: hardware not designed for multi-tenant, fault-tolerant operation creates utilization problems that no software layer can fully solve.""",

    9: """So here are the cloud-scale lessons for IC designers:

First, design for system deployability, not just benchmark performance. A chip that runs hot, draws excessive power, or lacks diagnostic capabilities creates downstream costs that dwarf the chip's purchase price.

Second, power efficiency is king. Performance per watt matters more than peak performance.

Third, build observability into the silicon. Remote monitoring, temperature reporting, error logging — at scale, they are operational necessities.

Fourth, reliability over years, not hours. We need chips that run stable at 85% load for 365 days, not chips that hit a record benchmark for 60 seconds.""",

    10: """Now I want to shift gears dramatically.

Everything I just described — liquid cooling, megawatt power supplies, high-speed interconnects — assumes the chip lives in a data center. But what happens when compute needs to leave the data center and walk around on two legs?

[Section transition — pause for effect]""",

    11: """A few years ago, I started asking: what comes after the data center? Where does AI compute go next? The answer is into the physical world — into robots, vehicles, and industrial equipment.

When I talk to the engineering teams, I always ask: what is your biggest bottleneck?

The answer is never funding. It's never algorithms. The answer is consistently: on-body compute. The available processors simply cannot deliver what we need within our power and size constraints.""",

    12: """LingShen Technology, or LivSyn — founded by a Tsinghua University team in 2023. They're building what they call a "universal brain platform" for humanoid robots.

What sets them apart is their four-in-one embodied large model that integrates perception, movement, manipulation, and planning. Their L1 robot has 26 degrees of freedom, weighs 85 kilograms, runs on NVIDIA Orin compute, and achieves over 12 hours of endurance.

What's remarkable from a deployment perspective: their pretrained models can generalize to a new scenario with fewer than 300 data samples, and be deployed in just 5 hours. They're already operating in smart manufacturing, hospital ward rounds, and data center inspection. Over 100 million RMB in Pre-A financing.""",

    13: """TIEN Kung — developed by the Beijing Innovation Center of Humanoid Robotics, a national-level platform supported by MIIT.

World's first purely electric-driven full-size humanoid robot capable of running at 6 km/h. Let me share some real numbers from their 2.0 Pro version: 173 centimeters tall, 73 kilograms, 42 degrees of freedom including dexterous hands, and on-body compute powered by dual NVIDIA Orin AGX modules delivering 550 TOPS — that's the kind of compute budget embodied AI has to work with.

Single-arm payload is 4 kilograms, bimanual up to 16. Battery life is 3.5 hours with hot-swappable battery packs. And critically, they've chosen an open-source platform strategy, making it a universal robot mother platform for the broader ecosystem.""",

    14: """Let me show you three real-world deployments that are already operating today.

First, smart manufacturing. LingShen's L1 robot is deployed on injection molding and stamping production lines — handling grasping, sorting, quality inspection, and assembly. Each robot replaces two to three manual workers, with ROI achieved within months.

Second, hospital ward rounds. The L1 robot assists doctors in completing morning rounds across 30 wards. It carries vital sign measurement devices — blood pressure, SpO2, temperature — syncs data to the hospital information system, and generates daily ward round reports automatically. It even conducts pre-diagnosis through medical large model and voice interaction.

Third, data center inspection. This one is close to my heart — 24/7 autonomous patrol of server racks. The robot detects anomalies like power failures, LED status changes, and can perform simple equipment restarts. And remember: it only needs 300 training samples and 5 hours to adapt to a new data center environment.

TIEN Kung's open platform is enabling applications in logistics, elder care, and hazardous environments.""",

    15: """Here is the comparison — and now I can give you real numbers from the robots I've invested in.

In the cloud, we measure in petaFLOPS. On TIEN Kung's body, total compute is 550 TOPS from dual Orin modules — that's all you get.

Power: 700 watts per accelerator in a liquid-cooled rack versus less than 50 watts total for on-body compute running on a 30-amp-hour battery. That's a 15x reduction.

Latency: one second is fine for cloud inference. For a robot catching its balance, you need under 10 milliseconds. Hesitate for one second and it falls.

Workload: in the cloud, typically one model at a time. On the robot body, you're running vision, large language model, motor control, and tactile processing simultaneously on the same silicon.

And cost: $30,000 per accelerator card is acceptable for data centers. But to reach mass adoption, on-body compute must be under $500 per unit.""",

    16: """[Pause — let this sink in]

Cloud computing's challenge is fitting chips into racks. Embodied computing's challenge is fitting racks into bodies.

These require fundamentally different IC design philosophies.""",

    17: """The good news: robots don't have to do everything locally. The emerging architecture is what I call the Big Brain, Mid Brain, Small Brain model.

The Big Brain is the cloud — our heterogeneous GPU pools running exaFLOPS of compute. This is where models are trained and world models are updated. It runs on the timescale of hours and days, and it's where innovations like heterogeneous resource pooling and vGPU slicing are maximizing utilization.

The Mid Brain is the edge — servers at the factory floor or warehouse. This is where multi-robot swarm coordination happens, connected via 5G or future 6G uplinks. Response time is in the hundreds of milliseconds.

The Small Brain is on-body compute — the SoC inside the robot itself. With 550 TOPS from dual Orin modules, it handles real-time perception, safety-critical control, and autonomous decision-making. Response time must be under 10 milliseconds.

[Walk through each tier, pointing to the diagram]""",

    18: """Each tier creates different IC demands.

Cloud tier — we know this well. GPU and AI accelerator market. Scale, power, interconnect.

Edge tier — where the inference chip market is exploding. High-throughput inference at 100-300 watts, multiple model types, real-time scheduling.

On-body tier — the wide-open frontier. Heterogeneous SoCs integrating CNN, transformer, and control processors, all under 50 watts.

This third tier is the single largest greenfield opportunity for IC designers in the next decade.""",

    19: """[Section transition]

Let me synthesize what I've learned from operating at cloud scale and investing at the robot edge into specific recommendations.

[Pause briefly]""",

    20: """For cloud and data center chip designers:

One — prioritize performance per watt. Your customers' biggest cost is electricity, not your chip.

Two — standardize interconnect interfaces. Proprietary interconnects create vendor lock-in that slows adoption.

Three — build diagnostics into the silicon. Remote health reporting, error counters, thermal sensors are not optional. They make large-scale operation viable.

Four — design for multi-tenant workloads. Hardware-level memory isolation, secure context switching, QoS guarantees directly improve utilization rates.""",

    21: """For edge and embodied chip designers:

One — heterogeneous compute is a requirement. CNN, transformer, and real-time control on a single die.

Two — power is THE hard constraint. Design to 50W first, then maximize compute. Sacrifice TOPS if needed.

Three — latency over throughput. One result in 5ms, not 1000 results per second.

Four — safety is non-negotiable. These chips operate inside machines alongside humans. Functional safety, error correction, graceful degradation — designed in from day one.

Five — cost under $500. A $200K robot will never reach mass adoption.""",

    22: """[Pause — key message]

The best chip is never the one with the highest benchmark score. It is the one that runs reliably, efficiently, and affordably in its target deployment environment.

This connects my entire career: from trading systems, through data centers, to robots. The principle is the same.""",

    23: """Looking forward:

2026 — the year of "inference famine." Demand for inference compute is growing faster than supply. Operators are racing to build out — some with plans for 50,000 petaFLOPS facilities. The gap will widen as embodied AI moves from labs to factories.

2027 — first wave of mass-produced humanoid robots in commercial deployment. Companies like LingShen already produce 25-30 units per month. Logistics, manufacturing, elder care — these are the first beachheads.

2028 to 2030 — the industry targets hardware costs dropping to around 30,000 RMB per robot, enabling C-end service scenarios. On-body AI compute market becomes comparable in scale to the smartphone SoC market. Every robot is a mobile computing platform.""",

    24: """[Convergence diagram]

Three markets converging: cloud training chips (mature), edge inference chips (growing), and on-body embodied SoCs (emerging).

The intersection is the IC opportunity — the core message of this talk.""",

    25: """For those of you in this room who design integrated circuits — you are designing the nervous system of tomorrow's intelligent machines.

Your decisions on architecture, power efficiency, and reliability will determine whether embodied AI becomes reality — or stays in the lab.

[Speak directly to the audience, make eye contact]""",

    26: """I started optimizing compute for financial models, where microseconds meant money.

I then built computing centers, where power efficiency meant survival.

Now I invest in robots, where real-time inference means safety.

At every stage, the lesson is the same: great chips don't just compute — they integrate into systems, they operate at scale, and they serve the real world.

The theme of this conference is "Advancing IC with AI." I'd like to add: we also need to advance AI with IC — from the cloud, through the edge, into the body.

[Pause for emphasis on last line]""",

    27: """Thank you.

[Wait for applause, then invite questions]

Q&A prepared topics:
- US chip export restrictions → focus on supply chain diversification driving innovation, avoid geopolitical commentary
- Domestic chips deployed → speak generally about maturing ecosystem, importance of deployment feedback loops. Mention heterogeneous pooling can compensate for individual chip performance gaps.
- Timeline for mass humanoid robots → 13,000 units shipped in 2025, expect 10x by 2028. LingShen already at 25-30/month. Industry targets 30K RMB cost by 2030. Bottleneck is compute, not mechanical hardware.
- Liquid cooling ROI → yes, unambiguously, for >30kW/rack, 18-24 month payback from energy savings alone. Waste heat recovery for agriculture is an emerging bonus.
- On-body compute specifics → TIEN Kung 2.0 Pro uses dual Orin AGX at 550 TOPS total, 42 DOF, 3.5h battery. LingShen L1 uses single Orin, 26 DOF, 12h+. Both NVIDIA platform — shows current dependence and opportunity for new entrants.
- 5G/6G for swarm robotics → emerging "mid brain" architecture uses dedicated wireless for multi-robot coordination, especially in industrial and logistics scenarios""",
}

for i, slide in enumerate(prs.slides):
    if i in NOTES:
        set_notes(slide, NOTES[i])

# ── Save ──
output_path = r"C:\Users\chenc\Claude Code\新闻公司\IEEE_ICTA_2026_Keynote_ChaoChen.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print(f"Notes added: {len(NOTES)} slides")
